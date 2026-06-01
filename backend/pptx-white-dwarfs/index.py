"""
Генерация PPTX-презентации по теме «Белые карлики: конец пути обычных звёзд»
(Строение Вселенной). 12 слайдов, академический стиль: тёмно-фиолетовый / белый / голубой.
Студент: Иванов Иван Иванович, группа АС-21.
"""

import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw


# ─── Цветовая палитра ────────────────────────────────────────────────────────
DEEP_PURPLE  = RGBColor(0x0E, 0x07, 0x2A)  # фон
MID_PURPLE   = RGBColor(0x1E, 0x0F, 0x4A)  # акцент
CYAN         = RGBColor(0x4D, 0xD0, 0xE1)  # голубой акцент
LIGHT_CYAN   = RGBColor(0xB2, 0xEB, 0xF2)  # светлый голубой
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xD6, 0xE8)
PALE_BLUE    = RGBColor(0xE0, 0xF7, 0xFA)
STAR_WHITE   = RGBColor(0xF0, 0xF4, 0xFF)
WARM_YELLOW  = RGBColor(0xFF, 0xEE, 0x88)


def make_space_bg(w=1920, h=1080, title_slide=False):
    img = Image.new("RGB", (w, h), (14, 7, 42))
    draw = ImageDraw.Draw(img)
    # Градиент — более светлый к центру
    for i in range(0, h, 6):
        factor = abs(i - h // 2) / (h // 2)
        r = int(14 + (1 - factor) * 8)
        g = int(7 + (1 - factor) * 5)
        b = int(42 + (1 - factor) * 20)
        draw.rectangle([0, i, w, i + 5], fill=(r, g, b))
    # Звёздное небо
    import random
    rnd = random.Random(42)
    for _ in range(200):
        x = rnd.randint(0, w)
        y = rnd.randint(0, h)
        br = rnd.randint(100, 255)
        sz = rnd.choice([1, 1, 1, 2])
        draw.ellipse([x, y, x + sz, y + sz], fill=(br, br, br))
    # Голубая полоска сверху
    draw.rectangle([0, 0, w, 6], fill=(77, 208, 225))
    draw.rectangle([0, 6, w, 10], fill=(0, 172, 193))
    # Голубая полоска снизу
    draw.rectangle([0, h - 10, w, h - 6], fill=(0, 172, 193))
    draw.rectangle([0, h - 6, w, h], fill=(77, 208, 225))
    if title_slide:
        draw.rectangle([80, h // 2 + 60, w - 80, h // 2 + 64], fill=(77, 208, 225))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_bg(slide, buf, prs):
    buf.seek(0)
    W, H = prs.slide_width, prs.slide_height
    pic = slide.shapes.add_picture(buf, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def txt(slide, text, left, top, width, height,
        size=28, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def multiline(slide, lines, left, top, width, height,
              sizes, bolds, colors, aligns, spaces=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aligns[i] if i < len(aligns) else PP_ALIGN.LEFT
        if spaces and i < len(spaces) and spaces[i]:
            p.space_before = Pt(spaces[i])
        run = p.add_run(); run.text = line
        run.font.size = Pt(sizes[i] if i < len(sizes) else 24)
        run.font.bold = bolds[i] if i < len(bolds) else False
        run.font.color.rgb = colors[i] if i < len(colors) else WHITE
        run.font.name = "Calibri"


def bullet_slide(prs, blank, title_text, bullets):
    s = prs.slides.add_slide(blank)
    set_bg(s, make_space_bg(), prs)
    bar = s.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()
    txt(s, title_text, Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=34, bold=True, color=CYAN)
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.13), Inches(5.4))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0: p.space_before = Pt(6)
        run = p.add_run()
        run.text = ("• " if not b.startswith("—") else "") + b
        run.font.size = Pt(22); run.font.color.rgb = LIGHT_GREY
        run.font.name = "Calibri"
    return s


def handler(event: dict, context) -> dict:
    """Генерирует PPTX 'Белые карлики: конец пути обычных звёзд' (Строение Вселенной)"""

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type",
        }, "body": ""}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    # ── Слайд 1: Титульный ────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, make_space_bg(title_slide=True), prs)

    txt(s1, "АСТРОНОМИЯ  •  Строение Вселенной",
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.6),
        size=18, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    multiline(s1,
        ["Белые карлики:", "конец пути обычных звёзд"],
        Inches(0.5), Inches(1.4), Inches(12.33), Inches(2.4),
        [54, 48], [True, False],
        [CYAN, WHITE],
        [PP_ALIGN.CENTER]*2, [0, 8])

    bar = s1.shapes.add_shape(1, Inches(3.5), Inches(4.15), Inches(6.33), Pt(2))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()

    multiline(s1,
        ["Выполнил: Иванов Иван Иванович", "Группа: АС-21"],
        Inches(0.5), Inches(4.35), Inches(12.33), Inches(1.0),
        [22, 22], [False, False], [LIGHT_GREY, LIGHT_GREY],
        [PP_ALIGN.CENTER]*2, [0, 6])

    txt(s1, "Исследовательская работа",
        Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5),
        size=16, color=CYAN, align=PP_ALIGN.CENTER, italic=True)

    # ── Слайд 2: Актуальность, цель, гипотеза ────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, make_space_bg(), prs)
    bar2 = s2.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = CYAN; bar2.line.fill.background()
    txt(s2, "Введение: цель, актуальность, гипотеза",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=32, bold=True, color=CYAN)

    box2 = s2.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.13), Inches(5.4))
    tf2 = box2.text_frame; tf2.word_wrap = True
    entries = [
        ("АКТУАЛЬНОСТЬ:",
         "Более 97% всех звёзд Млечного Пути закончат жизнь как белые карлики, включая наше Солнце. "
         "Изучение этих объектов позволяет понять финальные стадии звёздной эволюции, "
         "природу тёмной материи и происхождение химических элементов во Вселенной."),
        ("ЦЕЛЬ:",
         "Исследовать физические свойства белых карликов, механизмы их формирования "
         "и роль в строении и эволюции Вселенной на основе современных астрофизических данных."),
        ("ГИПОТЕЗА:",
         "Белые карлики являются не «мёртвыми» объектами, а активными участниками "
         "галактической химической эволюции: аккрецируя вещество, они способны вспыхивать "
         "как новые и сверхновые типа Ia, обогащая межзвёздную среду тяжёлыми элементами."),
    ]
    first = True
    for label, body in entries:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        if not first: p.space_before = Pt(10)
        run = p.add_run(); run.text = label
        run.font.size = Pt(20); run.font.bold = True
        run.font.color.rgb = CYAN; run.font.name = "Calibri"
        p2 = tf2.add_paragraph(); p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(18); r2.font.color.rgb = LIGHT_GREY; r2.font.name = "Calibri"

    # ── Слайд 3: Эволюция звёзд — путь к белому карлику ─────────────────────
    bullet_slide(prs, blank, "Путь звезды к белому карлику", [
        "Главная последовательность: звезда сжигает водород в гелий миллиарды лет.",
        "Красный гигант: водород в ядре иссякает, звезда расширяется в 100–200 раз.",
        "Планетарная туманность: оболочка сбрасывается — красивейшее зрелище во Вселенной.",
        "Белый карлик: остаётся горячее углеродно-кислородное ядро размером с Землю.",
        "Путь возможен для звёзд с массой до 8 масс Солнца (≈97% всех звёзд).",
        "Наше Солнце станет белым карликом через ~5 миллиардов лет.",
        "Ядро белого карлика поддерживается давлением вырожденного электронного газа.",
    ])

    # ── Слайд 4: Физические свойства белых карликов ──────────────────────────
    bullet_slide(prs, blank, "Физические свойства белых карликов", [
        "Масса: как у Солнца (0,5–1,4 M☉), но сжата до размеров Земли (~12 000 км).",
        "Плотность: 10⁶ г/см³ — ложка вещества весит ~5 тонн.",
        "Поверхностная температура: от 8 000 до 150 000 К (у молодых).",
        "Предел Чандрасекара: максимальная масса белого карлика — 1,4 M☉.",
        "Поверхностное притяжение: в 350 000 раз сильнее, чем у Земли.",
        "Состав: углерод и кислород (ядро), тонкая оболочка из гелия и водорода.",
        "Магнитное поле: у некоторых достигает 10⁹ Гс — сильнейшее среди звёздных остатков.",
    ])

    # ── Слайд 5: Классификация белых карликов ────────────────────────────────
    bullet_slide(prs, blank, "Классификация белых карликов", [
        "DA — водородная атмосфера (≈75% всех белых карликов), чисто белые.",
        "DB — гелиевая атмосфера (~16%), горячее 45 000 К.",
        "DC — нет спектральных линий, очень холодные (< 5 000 К).",
        "DQ — углеродные молекулы в атмосфере.",
        "DZ — металлические линии (загрязнение от поглощённых астероидов).",
        "DAP / DBP — магнитные белые карлики с поляризованным излучением.",
        "Пекулярные: некоторые имеют атмосферу с углеродом, водородом и гелием одновременно.",
    ])

    # ── Слайд 6: Белые карлики и сверхновые типа Ia ──────────────────────────
    bullet_slide(prs, blank, "Белые карлики и сверхновые типа Ia", [
        "В двойных системах белый карлик способен аккрецировать вещество со звезды-компаньона.",
        "При достижении предела Чандрасекара (1,4 M☉) — термоядерный взрыв: сверхновая Ia.",
        "Сверхновые Ia — «стандартные свечи»: одинаковая светимость → измерение расстояний.",
        "Открытие тёмной энергии (1998) основано на наблюдениях сверхновых Ia (Нобель-2011).",
        "Сценарий «двойного карлика»: слияние двух белых карликов → аналогичный взрыв.",
        "После взрыва белый карлик полностью разрушается — нейтронная звезда не образуется.",
        "За одну вспышку выбрасывается ≈1 M☉ вещества, богатого железом и никелем.",
    ])

    # ── Слайд 7: Охлаждение и «смерть» белого карлика ────────────────────────
    bullet_slide(prs, blank, "Охлаждение и финальная судьба", [
        "Белый карлик медленно остывает — нет источников энергии, только запасённое тепло.",
        "Теоретический предел: через ~10¹⁵ лет остынет до температуры фона (~3 К).",
        "Тогда он станет «чёрным карликом» — холодным невидимым объектом из вырожденного вещества.",
        "Чёрных карликов во Вселенной пока нет: её возраст 13,8 млрд лет — слишком мало.",
        "Самые холодные из известных белых карликов: ~2 700 К (WD J2147-4035).",
        "Кристаллизация ядра: углерод и кислород образуют твёрдую кристаллическую решётку.",
        "Открытие «кристаллических» белых карликов подтверждено данными Gaia (2019).",
    ])

    # ── Слайд 8: Роль в строении Вселенной ───────────────────────────────────
    bullet_slide(prs, blank, "Роль белых карликов в строении Вселенной", [
        "«Химическая фабрика»: планетарные туманности рассеивают C, N, O — строительные блоки жизни.",
        "Стандартные свечи: сверхновые Ia позволили открыть ускоряющееся расширение Вселенной.",
        "Кандидаты на роль MACHO: часть тёмной материи может состоять из холодных карликов.",
        "Источник информации об эволюции Галактики: белые карлики — «летопись» её истории.",
        "Загрязнённые белые карлики (DZ) — свидетельства поглощения планетоподобных тел.",
        "Миллисекундные пульсары могут образовываться при аккреции на нейтронные звезды рядом с карликами.",
        "Более 300 000 белых карликов каталогизировано в обзорах SDSS и Gaia.",
    ])

    # ── Слайд 9: Знаменитые белые карлики ────────────────────────────────────
    bullet_slide(prs, blank, "Знаменитые белые карлики", [
        "Сириус B: первый открытый белый карлик (1862), спутник ярчайшей звезды неба.",
        "40 Эридана B: открыт в 1783 г., ближайший к Земле из ярких белых карликов.",
        "BPM 37093 (Lucy): возможно, кристаллизованное ядро — «алмазная звезда» 10⁳⁴ карат.",
        "WD J2147-4035: один из самых старых и холодных (3 900 К), возраст ~10,7 млрд лет.",
        "GJ 440: белый карлик на расстоянии 15 св. лет, Т ≈ 9 200 К.",
        "SDSS J1228+1040: поглощает астероидоподобное вещество — диск из обломков.",
        "RE J0317-853: самый быстро вращающийся (~725 об/с), сильное магнитное поле.",
    ])

    # ── Слайд 10: Вывод — 5 ограничений ─────────────────────────────────────
    s10 = prs.slides.add_slide(blank)
    set_bg(s10, make_space_bg(), prs)
    bar10 = s10.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3))
    bar10.fill.solid(); bar10.fill.fore_color.rgb = CYAN; bar10.line.fill.background()
    txt(s10, "Вывод: что невозможно (5 пунктов)",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=32, bold=True, color=CYAN)

    box10 = s10.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.13), Inches(5.4))
    tf10 = box10.text_frame; tf10.word_wrap = True
    negatives = [
        "1. Невозможно наблюдать реальное охлаждение белого карлика в реальном времени: процесс занимает миллиарды лет.",
        "2. Невозможно напрямую измерить внутреннюю структуру белого карлика: он непрозрачен для обычного излучения.",
        "3. Невозможно предотвратить взрыв сверхновой Ia, если карлик в двойной системе превысит предел Чандрасекара.",
        "4. Невозможно существование белого карлика тяжелее 1,4 M☉: он неизбежно взорвётся или схлопнется.",
        "5. Невозможно пока обнаружить «чёрные карлики»: их возраст должен быть больше нынешнего возраста Вселенной.",
    ]
    for i, line in enumerate(negatives):
        p = tf10.paragraphs[0] if i == 0 else tf10.add_paragraph()
        if i > 0: p.space_before = Pt(8)
        run = p.add_run(); run.text = line
        run.font.size = Pt(20); run.font.color.rgb = RGBColor(0xFF, 0xAA, 0xAA)
        run.font.name = "Calibri"

    # ── Слайд 11: Вывод — 5 реальных возможностей ────────────────────────────
    s11 = prs.slides.add_slide(blank)
    set_bg(s11, make_space_bg(), prs)
    bar11 = s11.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3))
    bar11.fill.solid(); bar11.fill.fore_color.rgb = CYAN; bar11.line.fill.background()
    txt(s11, "Вывод: реальные возможности (5 пунктов)",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=32, bold=True, color=CYAN)

    box11 = s11.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.13), Inches(5.4))
    tf11 = box11.text_frame; tf11.word_wrap = True
    positives = [
        "1. Реально: использовать белые карлики как «стандартные свечи» для измерения масштабов Вселенной.",
        "2. Реально: изучать химический состав поглощённых планетоподобных тел через спектр загрязнённых карликов.",
        "3. Реально: применять астросейсмологию для исследования внутреннего строения через колебания поверхности.",
        "4. Реально: использовать данные о кристаллизации ядер для уточнения возраста звёздных скоплений и Галактики.",
        "5. Реально: обнаруживать планеты вокруг белых карликов — первые подтверждения уже получены (2020).",
    ]
    for i, line in enumerate(positives):
        p = tf11.paragraphs[0] if i == 0 else tf11.add_paragraph()
        if i > 0: p.space_before = Pt(8)
        run = p.add_run(); run.text = line
        run.font.size = Pt(20); run.font.color.rgb = RGBColor(0xAA, 0xFF, 0xAA)
        run.font.name = "Calibri"

    # ── Слайд 12: Приложение — 7 вопросов ────────────────────────────────────
    s12 = prs.slides.add_slide(blank)
    set_bg(s12, make_space_bg(), prs)
    bar12 = s12.shapes.add_shape(1, Inches(0.5), Inches(1.55), Inches(12.33), Pt(3))
    bar12.fill.solid(); bar12.fill.fore_color.rgb = CYAN; bar12.line.fill.background()
    txt(s12, "Приложение: контрольные вопросы",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=32, bold=True, color=CYAN)

    box12 = s12.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(12.13), Inches(5.4))
    tf12 = box12.text_frame; tf12.word_wrap = True
    questions = [
        "1. Какие звёзды могут стать белыми карликами и почему не все?",
        "2. Что такое предел Чандрасекара и каково его числовое значение?",
        "3. Почему вещество белого карлика обладает столь высокой плотностью?",
        "4. Как белые карлики связаны с открытием тёмной энергии?",
        "5. Что произойдёт с белым карликом, если он начнёт аккрецировать вещество от звезды-компаньона?",
        "6. Что такое «чёрный карлик» и почему таких объектов пока не существует?",
        "7. Как планетарные туманности, предшествующие белым карликам, влияют на химический состав Вселенной?",
    ]
    for i, q in enumerate(questions):
        p = tf12.paragraphs[0] if i == 0 else tf12.add_paragraph()
        if i > 0: p.space_before = Pt(8)
        run = p.add_run(); run.text = q
        run.font.size = Pt(20); run.font.color.rgb = LIGHT_CYAN
        run.font.name = "Calibri"

    # ── Сохраняем ─────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    data = base64.b64encode(buf.read()).decode("utf-8")

    return {
        "statusCode": 200,
        "headers": {"Access-Control-Allow-Origin": "*", "Content-Type": "application/json"},
        "body": '{"filename":"white_dwarfs.pptx","data":"' + data + '"}',
    }
