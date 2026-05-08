"""
Генерация PPTX-презентации 'Техника экологичного выставления границ'
в природном зелёно-бежевом стиле. Слайды 2-5 — пионы без изменений.
Остальные слайды — бежевый фон, зелёный текст, декоративные элементы.
Возвращает файл в base64 для скачивания.
"""

import base64
import io
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw


DARK_GREEN   = RGBColor(0x2D, 0x50, 0x16)
MID_GREEN    = RGBColor(0x4A, 0x7C, 0x2F)
LIGHT_GREEN  = RGBColor(0x6B, 0x9E, 0x3A)
BEIGE_BG     = RGBColor(0xF5, 0xF0, 0xE8)
WARM_BEIGE   = RGBColor(0xE8, 0xDF, 0xC8)
DARK_BROWN   = RGBColor(0x3D, 0x2B, 0x1F)
RED_ACCENT   = RGBColor(0xCC, 0x22, 0x22)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

LOGO_URL    = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/29d08a49-613a-4d33-8137-21a389ad93de.jpg"

SLIDE2_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/00ee0148-ebfd-4675-84b3-d75532822211.jpg"
SLIDE3_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/61853eaf-fd3f-4c91-9f7a-f258abeee6ec.jpg"
SLIDE4_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/7110ae58-7bd6-4ecf-a2dc-23612fc7facd.jpg"
SLIDE5_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/208ebcf1-8745-4efd-99cc-633c2014858d.jpg"


def download_image_as_jpg(url: str, max_dim: int = 1920, bg_color=(245, 240, 232)) -> io.BytesIO:
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=25) as r:
        raw = io.BytesIO(r.read())
    img = Image.open(raw)
    if img.mode in ("RGBA", "P", "LA"):
        bg = Image.new("RGB", img.size, bg_color)
        if img.mode == "P":
            img = img.convert("RGBA")
        mask = img.split()[-1] if img.mode in ("RGBA", "LA") else None
        bg.paste(img.convert("RGBA") if img.mode not in ("RGBA",) else img, mask=mask)
        img = bg
    elif img.mode != "RGB":
        img = img.convert("RGB")
    w, h = img.size
    if max(w, h) > max_dim:
        ratio = max_dim / max(w, h)
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85, optimize=True)
    buf.seek(0)
    return buf


def make_beige_bg(width_px: int = 1920, height_px: int = 1080) -> io.BytesIO:
    img = Image.new("RGB", (width_px, height_px), (245, 240, 232))
    draw = ImageDraw.Draw(img)
    # Зелёная полоска сверху
    draw.rectangle([0, 0, width_px, 18], fill=(74, 124, 47))
    # Зелёная полоска снизу
    draw.rectangle([0, height_px - 18, width_px, height_px], fill=(74, 124, 47))
    # Тонкая линия под верхней полоской
    draw.rectangle([0, 18, width_px, 24], fill=(107, 158, 58))
    # Тонкая линия над нижней полоской
    draw.rectangle([0, height_px - 24, width_px, height_px - 18], fill=(107, 158, 58))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def make_title_bg(width_px: int = 1920, height_px: int = 1080) -> io.BytesIO:
    """Фон для титульного слайда — бежевый с зелёным боковым акцентом"""
    img = Image.new("RGB", (width_px, height_px), (245, 240, 232))
    draw = ImageDraw.Draw(img)
    # Левая зелёная полоса
    draw.rectangle([0, 0, 12, height_px], fill=(45, 80, 22))
    draw.rectangle([12, 0, 22, height_px], fill=(74, 124, 47))
    # Правая зелёная полоса
    draw.rectangle([width_px - 22, 0, width_px - 12, height_px], fill=(74, 124, 47))
    draw.rectangle([width_px - 12, 0, width_px, height_px], fill=(45, 80, 22))
    # Верхняя полоска
    draw.rectangle([0, 0, width_px, 10], fill=(45, 80, 22))
    # Нижняя полоска
    draw.rectangle([0, height_px - 10, width_px, height_px], fill=(45, 80, 22))
    # Декоративный прямоугольник в центре внизу (бежевый тёплый)
    draw.rectangle([22, height_px - 160, width_px - 22, height_px - 10], fill=(232, 223, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_bg(slide, img_buf: io.BytesIO, prs):
    img_buf.seek(0)
    W = prs.slide_width
    H = prs.slide_height
    pic = slide.shapes.add_picture(img_buf, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_full_image_slide(prs, slide, url):
    try:
        img_buf = download_image_as_jpg(url)
        W = prs.slide_width
        H = prs.slide_height
        img_buf.seek(0)
        pil_img = Image.open(img_buf)
        orig_w, orig_h = pil_img.size
        scale = max(W / orig_w, H / orig_h)
        new_w = int(orig_w * scale)
        new_h = int(orig_h * scale)
        offset_x = (W - new_w) // 2
        offset_y = (H - new_h) // 2
        img_buf.seek(0)
        pic = slide.shapes.add_picture(img_buf, offset_x, offset_y, new_w, new_h)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    except Exception:
        pass


def add_text(slide, text, left, top, width, height,
             font_size=40, bold=False, color=DARK_GREEN,
             align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Georgia"
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_sizes, bolds, colors, aligns, space_before_pts=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aligns[i] if i < len(aligns) else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_sizes[i] if i < len(font_sizes) else 40)
        run.font.bold = bolds[i] if i < len(bolds) else False
        run.font.color.rgb = colors[i] if i < len(colors) else DARK_GREEN
        run.font.name = "Georgia"
        if i > 0:
            sp = space_before_pts[i] if space_before_pts and i < len(space_before_pts) else 14
            p.space_before = Pt(sp)
    return txBox


def handler(event: dict, context) -> dict:
    """Генерирует PPTX-презентацию 'Техника экологичного выставления границ' в природном зелёно-бежевом стиле"""

    if event.get("httpMethod") == "OPTIONS":
        return {
            "statusCode": 200,
            "headers": {
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
                "Access-Control-Allow-Headers": "Content-Type",
                "Access-Control-Max-Age": "86400",
            },
            "body": "",
        }

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    # ── Слайд 1: Титульный ───────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank)
    set_bg(slide1, make_title_bg(), prs)

    # Заголовок
    add_text(slide1, "Техника экологичного\nвыставления границ",
             Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.2),
             font_size=52, bold=False, color=DARK_GREEN, align=PP_ALIGN.CENTER)

    # Разделитель — зелёная линия
    sep = slide1.shapes.add_shape(1, Inches(4.5), Inches(4.2), Inches(4.3), Pt(3))
    sep.fill.solid()
    sep.fill.fore_color.rgb = MID_GREEN
    sep.line.fill.background()

    # Подпись центра
    add_text(slide1, "Центр квантовой педагогики и психологии «Фуллерен»",
             Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.7),
             font_size=22, bold=False, color=MID_GREEN, align=PP_ALIGN.CENTER)

    # ── Слайды 2-5: Пионы без изменений ─────────────────────────────────────
    for url in [SLIDE2_URL, SLIDE3_URL, SLIDE4_URL, SLIDE5_URL]:
        slide = prs.slides.add_slide(blank)
        add_full_image_slide(prs, slide, url)

    # ── Вспомогательная функция для текстового слайда ────────────────────────
    def text_slide(lines, font_sizes, bolds, colors, space_pts=None):
        s = prs.slides.add_slide(blank)
        set_bg(s, make_beige_bg(), prs)
        aligns = [PP_ALIGN.CENTER] * len(lines)
        total_lines = len(lines)
        block_h = Inches(1.0) * total_lines + Inches(0.5)
        top = (H - block_h) / 2 - Inches(0.3)
        add_multiline(s, lines,
                      Inches(1.0), top, Inches(11.33), block_h,
                      font_sizes, bolds, colors, aligns, space_pts)
        return s

    # ── Слайд 6: ГОРОД ───────────────────────────────────────────────────────
    text_slide(["ГОРОД"], [64], [True], [DARK_GREEN])

    # ── Слайд 7: ГОРОД / ГРАД ────────────────────────────────────────────────
    text_slide(["ГОРОД", "ГРАД"], [54, 54], [True, True], [DARK_GREEN, DARK_GREEN], [0, 20])

    # ── Слайд 8: ГОРОД / ГРАД / ОГРАДА ───────────────────────────────────────
    text_slide(["ГОРОД", "ГРАД", "ОГРАДА"], [48, 48, 48], [True, True, True],
               [DARK_GREEN, DARK_GREEN, MID_GREEN], [0, 18, 18])

    # ── Слайд 9: КТО БЫЛ ПЕРВЫМ ГРАДОСТРОИТЕЛЕМ? ─────────────────────────────
    text_slide(["КТО БЫЛ ПЕРВЫМ\nГРАДОСТРОИТЕЛЕМ?"], [52], [True], [DARK_GREEN])

    # ── Слайд 10: Ты будешь сКИТАльцем по всей Земле. ────────────────────────
    text_slide(["Ты будешь\nсКИТАльцем\nпо всей Земле."], [48], [False], [DARK_BROWN])

    # ── Слайд 11: Корень ограничений / СТРАХ ─────────────────────────────────
    text_slide(["Корень ограничений", "СТРАХ"],
               [44, 60], [False, True],
               [DARK_GREEN, RGBColor(0x8B, 0x1A, 0x1A)], [0, 24])

    # ── Слайд 12: БЕЗГРАНИЧНОСТЬ ─────────────────────────────────────────────
    text_slide(["БЕЗГРАНИЧНОСТЬ"], [64], [True], [DARK_GREEN])

    # ── Слайд 13: БЕЗГРАНИЧНОСТЬ / по ВЕРе каждому дано будет ────────────────
    text_slide(["БЕЗГРАНИЧНОСТЬ", "по ВЕРе каждому дано будет"],
               [52, 44], [True, False],
               [DARK_GREEN, DARK_GREEN], [0, 28])

    # ── Слайд 14: БЕЗГРАНИЧНОСТЬ / по ВЕРе... / уВЕРенность ─────────────────
    text_slide(["БЕЗГРАНИЧНОСТЬ", "по ВЕРе каждому дано будет", "уВЕРенность"],
               [44, 38, 48], [True, False, True],
               [DARK_GREEN, DARK_GREEN, MID_GREEN], [0, 24, 24])

    # ── Слайд 15: Как проявляется уВЕРенность? ───────────────────────────────
    text_slide(["Как проявляется\nуВЕРенность?"], [52], [False], [DARK_GREEN])

    # ── Слайд 16: не нравится / отстаньте / не хочу ──────────────────────────
    text_slide(["не нравится", "отстаньте", "не хочу"],
               [48, 48, 48], [False, False, False],
               [DARK_BROWN, DARK_BROWN, DARK_BROWN], [0, 16, 16])

    # ── Слайд 17: не нравится / отстаньте / не хочу / СТОП!!! ───────────────
    text_slide(["не нравится", "отстаньте", "не хочу", "СТОП!!!"],
               [44, 44, 44, 58], [False, False, False, True],
               [DARK_BROWN, DARK_BROWN, DARK_BROWN, RED_ACCENT], [0, 14, 14, 28])

    # ── Слайд 18: ВМЕШАТЕЛЬСТВО / КАК ТРЕТЬЯ СТОРОНА / В КОНФЛИКТ ───────────
    text_slide(["ВМЕШАТЕЛЬСТВО", "КАК ТРЕТЬЯ СТОРОНА", "В КОНФЛИКТ"],
               [54, 46, 54], [True, False, True],
               [DARK_GREEN, MID_GREEN, DARK_GREEN], [0, 28, 28])

    # ── Слайд 19: КАК МЫ МОЖЕМ НАРУШАТЬ ГРАНИЦЫ? ────────────────────────────
    text_slide(["КАК МЫ МОЖЕМ\nНАРУШАТЬ\nГРАНИЦЫ?"], [52], [True], [DARK_GREEN])

    # ── Слайд 20: СОВЕТЫ (ПРЕДЛОЖЕНИЯ) / 5 ВИДОВ ВМЕШАТЕЛЬСТВ ───────────────
    text_slide(["СОВЕТЫ (ПРЕДЛОЖЕНИЯ)", "5 ВИДОВ ВМЕШАТЕЛЬСТВ"],
               [48, 56], [False, True],
               [MID_GREEN, DARK_GREEN], [0, 32])

    # ── Сохраняем и отдаём base64 ─────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    data = base64.b64encode(buf.read()).decode("utf-8")

    return {
        "statusCode": 200,
        "headers": {
            "Access-Control-Allow-Origin": "*",
            "Content-Type": "application/json",
        },
        "body": '{"filename":"granitsy_presentation.pptx","data":"' + data + '"}',
    }
