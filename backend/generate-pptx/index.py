"""
Генерация PPTX-презентации "Квантовая педагогика" в золотисто-тёмном стиле.
Возвращает файл в base64 для скачивания.
"""

import base64
import io
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw


# ── Цвета ────────────────────────────────────────────────────────────────────
GOLD      = RGBColor(0xFF, 0xD7, 0x00)
GOLD2     = RGBColor(0xFF, 0xA5, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF0, 0xE8, 0xD0)
DARK_GOLD = RGBColor(0x3A, 0x28, 0x00)
DEEP_DARK = RGBColor(0x06, 0x04, 0x00)

# ── URL картинок ──────────────────────────────────────────────────────────────
SPIRAL_URL   = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/4083467d-4f30-4338-bbe4-204644e9c2bf.png"
BAZARNY_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/2b075abe-8b7b-4015-bfb3-71ee7198df9a.png"
SHCHET_URL   = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/1589e3bb-bdc7-4e1e-9977-28bab6bc26a9.png"
AMON_URL     = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/c70ffe39-d7ee-4e81-87cc-89864954889a.png"
SHAT_URL     = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/8aa076ab-aa1c-43c5-b471-b901927f46f7.png"
SOKRAT_URL   = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/c8ffe4c7-d47b-4235-b48c-7a644fc570ce.png"
APPLE_URL    = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/05cdfdea-ae9b-40b3-8681-1e4bb6e6cbd2.png"
TREE_URL     = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/ff7342ac-235a-4002-a51a-befd9f5c20f7.png"
EARTH_URL    = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/7879c67d-7b9d-4ba0-aa69-1e9f4d2659f4.png"
BORSCH_URL   = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/a66ac4ea-a927-43d4-a9fd-d3658863892c.png"
QUANTUM_URL  = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/e045064e-c54e-4f2d-a5d9-60c9cd417c0e.jpg"
EQUATION_URL = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/e726d71e-7965-4cf4-971e-0c2022653dc8.png"
CIRCLE_URL   = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/2f8895f9-2661-4dc9-9754-242172465263.png"
TITLE_BG_URL = "https://cdn.poehali.dev/projects/a853d61a-73f8-407d-846b-967c4543637c/bucket/fa0cedf0-1410-4c4e-9166-7d741ba6be9a.png"


def download_image(url: str, max_dim: int = 800) -> io.BytesIO:
    """Скачивает, масштабирует и конвертирует в JPEG для уменьшения размера"""
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=25) as r:
        raw = io.BytesIO(r.read())
    img = Image.open(raw)
    if img.mode in ("RGBA", "P", "LA"):
        bg = Image.new("RGB", img.size, (10, 8, 2))
        if img.mode == "P":
            img = img.convert("RGBA")
        bg.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
        img = bg
    elif img.mode != "RGB":
        img = img.convert("RGB")
    w, h = img.size
    if max(w, h) > max_dim:
        ratio = max_dim / max(w, h)
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=72, optimize=True)
    buf.seek(0)
    return buf


def make_photo_bg(width_px: int, height_px: int) -> io.BytesIO:
    img = Image.new("RGB", (width_px, height_px), (8, 6, 0))
    draw = ImageDraw.Draw(img)
    for y in range(height_px // 2, height_px):
        ratio = (y - height_px // 2) / (height_px // 2)
        r = int(8 + ratio * 60)
        g = int(6 + ratio * 40)
        for x in range(width_px):
            draw.point((x, y), fill=(r, g, 0))
    for x in range(width_px * 3 // 4, width_px):
        ratio = (x - width_px * 3 // 4) / (width_px // 4)
        r = int(50 * ratio)
        g = int(30 * ratio)
        for y in range(height_px):
            px = img.getpixel((x, y))
            draw.point((x, y), fill=(min(255, px[0] + r), min(255, px[1] + g), 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def make_dark_bg(width_px: int, height_px: int) -> io.BytesIO:
    img = Image.new("RGB", (width_px, height_px), (10, 8, 2))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_bg_image(slide, img_buf: io.BytesIO, prs):
    img_buf.seek(0)
    W = prs.slide_width
    H = prs.slide_height
    pic = slide.shapes.add_picture(img_buf, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_photo_proportional(slide, img_buf, left, top, max_w, max_h):
    img_buf.seek(0)
    pil_img = Image.open(img_buf)
    orig_w, orig_h = pil_img.size
    ratio = min(max_w / orig_w, max_h / orig_h)
    new_w = int(orig_w * ratio)
    new_h = int(orig_h * ratio)
    offset_x = (max_w - new_w) // 2
    offset_y = (max_h - new_h) // 2
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left + offset_x, top + offset_y, new_w, new_h)


def build_image_text_slide(slide, prs, img_raw, text,
                            text_color=WHITE, font_size=28, bold=False, italic=False,
                            img_right=True, bg_buf=None):
    W = prs.slide_width
    H = prs.slide_height
    if bg_buf:
        bg_buf.seek(0)
        set_bg_image(slide, bg_buf, prs)
    add_rect(slide, 0, 0, W, Inches(0.45), DARK_GOLD)
    half = W // 2
    if img_right:
        add_textbox(slide, text,
                    Inches(0.5), Inches(0.9), half - Inches(0.6), H - Inches(1.5),
                    font_size=font_size, bold=bold, italic=italic,
                    color=text_color, align=PP_ALIGN.LEFT)
        if img_raw:
            img_raw.seek(0)
            add_photo_proportional(slide, img_raw,
                                   half + Inches(0.2), Inches(0.55),
                                   half - Inches(0.4), H - Inches(0.7))
    else:
        if img_raw:
            img_raw.seek(0)
            add_photo_proportional(slide, img_raw,
                                   Inches(0.2), Inches(0.55),
                                   half - Inches(0.4), H - Inches(0.7))
        add_textbox(slide, text,
                    half + Inches(0.3), Inches(0.9), half - Inches(0.6), H - Inches(1.5),
                    font_size=font_size, bold=bold, italic=italic,
                    color=text_color, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), DARK_GOLD)


def handler(event: dict, context) -> dict:
    """Генерирует PPTX-презентацию о квантовой педагогике (золотой стиль, 16 слайдов)"""

    if event.get("httpMethod") == "OPTIONS":
        return {
            "statusCode": 200,
            "headers": {
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
                "Access-Control-Allow-Headers": "Content-Type",
                "Access-Control-Max-Age": "86400",
            },
            "body": "",
        }

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    def safe_dl(url):
        try:
            return download_image(url)
        except:
            return None

    # Скачиваем все картинки
    title_bg_raw = safe_dl(TITLE_BG_URL)
    spiral_raw   = safe_dl(SPIRAL_URL)
    bazarny_raw  = safe_dl(BAZARNY_URL)
    shchet_raw   = safe_dl(SHCHET_URL)
    amon_raw     = safe_dl(AMON_URL)
    shat_raw     = safe_dl(SHAT_URL)
    sokrat_raw   = safe_dl(SOKRAT_URL)
    apple_raw    = safe_dl(APPLE_URL)
    tree_raw     = safe_dl(TREE_URL)
    earth_raw    = safe_dl(EARTH_URL)
    borsch_raw   = safe_dl(BORSCH_URL)
    quantum_raw  = safe_dl(QUANTUM_URL)
    equation_raw = safe_dl(EQUATION_URL)
    circle_raw   = safe_dl(CIRCLE_URL)

    bg_photo = make_photo_bg(1333, 750)
    bg_dark  = make_dark_bg(1333, 750)

    photos_3 = [
        (shchet_raw,  "Михаил Петрович\nЩетинин"),
        (amon_raw,    "Шалва Александрович\nАмонашвили"),
        (shat_raw,    "Виктор Фёдорович\nШаталов"),
    ]

    # ── СЛАЙД 1: Титульный — просто картинка на весь слайд ───────────────────
    s1 = prs.slides.add_slide(blank)
    if title_bg_raw:
        title_bg_raw.seek(0)
        set_bg_image(s1, title_bg_raw, prs)

    # ── СЛАЙД 2: Базарный ────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    bg_photo.seek(0)
    set_bg_image(s2, bg_photo, prs)
    add_rect(s2, 0, 0, W, Inches(0.55), DARK_GOLD)
    add_textbox(s2, "Владимир Филиппович Базарный",
                Inches(0.4), Inches(0.05), W - Inches(0.8), Inches(0.45),
                font_size=26, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    if bazarny_raw:
        bazarny_raw.seek(0)
        add_photo_proportional(s2, bazarny_raw, Inches(3.5), Inches(0.7), Inches(6.5), Inches(6.2))
    add_rect(s2, 0, H - Inches(0.5), W, Inches(0.5), DARK_GOLD)

    # ── Вспомогательная: три фото ─────────────────────────────────────────────
    def build_three(slide_obj, bottom_text, bottom_italic=False, text_color=WHITE):
        bg_photo.seek(0)
        set_bg_image(slide_obj, bg_photo, prs)
        add_rect(slide_obj, 0, 0, W, Inches(0.45), DARK_GOLD)
        photo_max_w = Inches(3.9)
        photo_max_h = Inches(4.2)
        gap = (W - photo_max_w * 3) // 4
        for i, (raw, name) in enumerate(photos_3):
            left = gap + i * (photo_max_w + gap)
            if raw:
                raw.seek(0)
                add_photo_proportional(slide_obj, raw, left, Inches(0.55), photo_max_w, photo_max_h)
            add_textbox(slide_obj, name, left, Inches(4.85), photo_max_w, Inches(0.75),
                        font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        bar_h = Inches(1.5)
        add_rect(slide_obj, 0, H - bar_h, W, bar_h, DARK_GOLD)
        add_textbox(slide_obj, bottom_text,
                    Inches(0.5), H - bar_h + Inches(0.1), W - Inches(1), bar_h - Inches(0.15),
                    font_size=24, bold=not bottom_italic, italic=bottom_italic,
                    color=text_color, align=PP_ALIGN.CENTER)

    # ── СЛАЙДЫ 3-5 ───────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    build_three(s3,
        '"Почему мы не можем с таким же успехом повторить опыт этих личностей?"',
        bottom_italic=True)

    s4 = prs.slides.add_slide(blank)
    build_three(s4, "Учит не методика этих людей", text_color=OFF_WHITE)

    s5 = prs.slides.add_slide(blank)
    build_three(s5, "Сама ЛИЧНОСТЬ преОБРАЗует пространство", text_color=GOLD)

    # ── СЛАЙД 6: Сократ + вопрос ─────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s6, prs, sokrat_raw,
        "Где найти такого учителя,\nкоторый своим состоянием\n«Вдохновения»\nменяет пространство?",
        text_color=WHITE, font_size=28, bold=False, italic=True,
        img_right=False, bg_buf=bg_dark)

    # ── СЛАЙД 7: Сократ + ответ ───────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s7, prs, sokrat_raw,
        "Готов ученик —\nпридёт и учитель",
        text_color=GOLD, font_size=38, bold=True,
        img_right=False, bg_buf=bg_dark)

    # ── СЛАЙД 8: Яблоко + вопросы ─────────────────────────────────────────────
    s8 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s8, prs, apple_raw,
        "Кто ученик?\n\nКак подготовить ученика,\nчтобы к нему пришёл учитель?",
        text_color=WHITE, font_size=27, bold=False, italic=True,
        img_right=True, bg_buf=bg_dark)

    # ── СЛАЙД 9: Дерево с яблоком ────────────────────────────────────────────
    s9 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s9, prs, tree_raw,
        "Плод зависит от\nСОСТОЯНИЯ дерева",
        text_color=GOLD, font_size=36, bold=True,
        img_right=True, bg_buf=bg_dark)

    # ── СЛАЙД 10: Дерево с землёй ─────────────────────────────────────────────
    s10 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s10, prs, earth_raw,
        "Дерево зависит от\nСОСТОЯНИЯ почвы",
        text_color=GOLD, font_size=36, bold=True,
        img_right=True, bg_buf=bg_dark)

    # ── СЛАЙД 11: Борщ ────────────────────────────────────────────────────────
    s11 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s11, prs, borsch_raw,
        "слагаемые\nне меняются местами",
        text_color=WHITE, font_size=34, bold=True,
        img_right=False, bg_buf=bg_dark)

    # ── СЛАЙД 12: Дерево с землёй + длинный текст ────────────────────────────
    s12 = prs.slides.add_slide(blank)
    bg_dark.seek(0)
    build_image_text_slide(s12, prs, earth_raw,
        "Важна последовательность:\nсначала СОСТОЯНИЕ мамы,\nпотом СОСТОЯНИЕ папы,\nа затем преОБРАЗование плода",
        text_color=OFF_WHITE, font_size=25, bold=False,
        img_right=False, bg_buf=bg_dark)

    # ── СЛАЙДЫ 13-15: Квантовая запутанность ─────────────────────────────────
    bar_h = Inches(1.7)

    s13 = prs.slides.add_slide(blank)
    if quantum_raw:
        quantum_raw.seek(0)
        set_bg_image(s13, quantum_raw, prs)
    add_rect(s13, 0, H - bar_h, W, bar_h, DARK_GOLD)
    add_textbox(s13,
        "Переход от ОБРАЗования, как подлежащего,\nк ОБРАЗованию, как СКАЗУЕМОГО",
        Inches(0.5), H - bar_h + Inches(0.1), W - Inches(1), bar_h - Inches(0.15),
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    s14 = prs.slides.add_slide(blank)
    if quantum_raw:
        quantum_raw.seek(0)
        set_bg_image(s14, quantum_raw, prs)
    add_rect(s14, 0, H - bar_h, W, bar_h, DARK_GOLD)
    add_textbox(s14,
        "Только изменения самого СОСТОЯНИЯ родителя\nсможет преОБРАЗовать готовность ребёнка к обучению",
        Inches(0.5), H - bar_h + Inches(0.1), W - Inches(1), bar_h - Inches(0.15),
        font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    s15 = prs.slides.add_slide(blank)
    if quantum_raw:
        quantum_raw.seek(0)
        set_bg_image(s15, quantum_raw, prs)
    add_rect(s15, 0, H - Inches(1.3), W, Inches(1.3), DARK_GOLD)
    add_textbox(s15,
        "закон квантовой запутанности",
        Inches(0.5), H - Inches(1.25), W - Inches(1), Inches(1.2),
        font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── СЛАЙД 16: Уравнения + Круг + мотивация ───────────────────────────────
    s16 = prs.slides.add_slide(blank)
    if equation_raw:
        equation_raw.seek(0)
        set_bg_image(s16, equation_raw, prs)

    # Тёмная подложка слева
    add_rect(s16, 0, 0, Inches(6.6), H, RGBColor(0x05, 0x08, 0x18))
    add_rect(s16, 0, 0, Inches(6.6), Inches(0.55), DARK_GOLD)
    add_textbox(s16, "Мотивация",
                Inches(0.3), Inches(0.06), Inches(6.1), Inches(0.45),
                font_size=22, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    add_textbox(s16, "Внешняя мотивация:",
                Inches(0.3), Inches(0.7), Inches(6.0), Inches(0.5),
                font_size=20, bold=True, color=GOLD2, align=PP_ALIGN.LEFT)
    add_textbox(s16,
        "— сдать гос. аттестацию;\n— поступить в Университет;\n— найти хорошую работу;\n— приносить пользу.",
        Inches(0.3), Inches(1.25), Inches(6.0), Inches(2.0),
        font_size=17, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_textbox(s16, "Внутренняя мотивация:",
                Inches(0.3), Inches(3.4), Inches(6.0), Inches(0.5),
                font_size=20, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_textbox(s16,
        "— нужно ли мне разобраться со справедливостью в своей жизни?\n— как уравнения мне помогут это сделать?",
        Inches(0.3), Inches(3.95), Inches(6.0), Inches(2.2),
        font_size=17, bold=False, color=OFF_WHITE, align=PP_ALIGN.LEFT)

    add_rect(s16, 0, H - Inches(0.35), W, Inches(0.35), DARK_GOLD)



    # ── Сохраняем ────────────────────────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    pptx_b64 = base64.b64encode(output.read()).decode("utf-8")

    return {
        "statusCode": 200,
        "headers": {
            "Access-Control-Allow-Origin": "*",
            "Content-Type": "application/json",
        },
        "body": f'{{"filename": "kvantovaya_pedagogika.pptx", "data": "{pptx_b64}"}}',
    }