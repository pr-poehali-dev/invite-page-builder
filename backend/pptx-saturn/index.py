"""
Генерация PPTX-презентации по теме №22:
«Космические зонды к внешним планетам: Voyager и открытия колец Урана»
(Солнечная система). 13 слайдов, тёмно-синий / золотой стиль, с картинками.
Студент: Ефремова Анастасия Александровна, группа 11С.
"""

import base64
import io
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw


# ─── Цветовая палитра ────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x07, 0x10, 0x2E)
MID_BLUE    = RGBColor(0x12, 0x23, 0x55)
GOLD        = RGBColor(0xD4, 0xA9, 0x17)
LIGHT_GOLD  = RGBColor(0xF0, 0xD0, 0x6A)
CYAN        = RGBColor(0x4D, 0xD0, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xD6, 0xE8)
PALE_BLUE   = RGBColor(0xB2, 0xD8, 0xF0)
RED_SOFT    = RGBColor(0xFF, 0x8A, 0x80)
GREEN_SOFT  = RGBColor(0xA5, 0xD6, 0xA7)

# ─── Картинки из NASA/открытых источников ────────────────────────────────────
IMG_VOYAGER = "https://upload.wikimedia.org/wikipedia/commons/thumb/6/6d/Voyager_spacecraft.jpg/1280px-Voyager_spacecraft.jpg"
IMG_URANUS  = "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3d/Uranus2.jpg/1280px-Uranus2.jpg"
IMG_RINGS   = "https://upload.wikimedia.org/wikipedia/commons/thumb/1/1f/Uranian_rings_PIA01977.jpg/1280px-Uranian_rings_PIA01977.jpg"
IMG_NEPTUNE = "https://upload.wikimedia.org/wikipedia/commons/thumb/b/b9/Neptune_Voyager2_color_calibrated.png/1280px-Neptune_Voyager2_color_calibrated.png"
IMG_SOLAR   = "https://upload.wikimedia.org/wikipedia/commons/thumb/c/cb/Planets2013.svg/1280px-Planets2013.svg.png"


def fetch_image(url: str, max_dim: int = 480) -> io.BytesIO | None:
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=20) as r:
            raw = io.BytesIO(r.read())
        img = Image.open(raw)
        # thumbnail не увеличивает, только уменьшает — экономит память
        img.thumbnail((max_dim, max_dim), Image.LANCZOS)
        if img.mode in ("RGBA", "P", "LA"):
            bg = Image.new("RGB", img.size, (7, 16, 46))
            if img.mode == "P":
                img = img.convert("RGBA")
            alpha = img.split()[-1] if img.mode in ("RGBA", "LA") else None
            bg.paste(img.convert("RGBA"), mask=alpha)
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=72)
        buf.seek(0)
        return buf
    except Exception:
        return None


def make_bg(w=960, h=540, title=False) -> io.BytesIO:
    img = Image.new("RGB", (w, h), (7, 16, 46))
    draw = ImageDraw.Draw(img)
    # Тонкий вертикальный градиент
    for i in range(0, h, 3):
        f = i / h
        r = int(7 + f * 5)
        g = int(16 + f * 8)
        b = int(46 + f * 20)
        draw.rectangle([0, i, w, i + 2], fill=(r, g, b))
    # Звёзды
    import random
    rnd = random.Random(99)
    for _ in range(250):
        x = rnd.randint(0, w)
        y = rnd.randint(0, h)
        br = rnd.randint(90, 240)
        sz = rnd.choice([1, 1, 1, 2])
        draw.ellipse([x, y, x + sz, y + sz], fill=(br, br, br))
    # Полоски
    draw.rectangle([0, 0, w, 7], fill=(212, 169, 23))
    draw.rectangle([0, 7, w, 11], fill=(150, 110, 10))
    draw.rectangle([0, h - 11, w, h - 7], fill=(150, 110, 10))
    draw.rectangle([0, h - 7, w, h], fill=(212, 169, 23))
    if title:
        # Декоративная горизонтальная линия для титула
        draw.rectangle([100, h // 2 + 80, w - 100, h // 2 + 83], fill=(212, 169, 23))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_bg(slide, buf, prs):
    buf.seek(0)
    W, H = prs.slide_width, prs.slide_height
    pic = slide.shapes.add_picture(buf, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_image_right(slide, prs, url: str,
                    left_frac=0.57, top_frac=0.18,
                    w_frac=0.38, h_frac=0.70,
                    round_corners=True):
    """Добавляет картинку справа — простой вариант без тяжёлого resize."""
    W, H = prs.slide_width, prs.slide_height
    buf = fetch_image(url)
    if buf is None:
        return
    left = int(W * left_frac)
    top  = int(H * top_frac)
    w    = int(W * w_frac)
    h    = int(H * h_frac)
    # Просто вставляем уже уменьшенную картинку, pptx растянет по w/h
    slide.shapes.add_picture(buf, left, top, w, h)


def txt(slide, text, left, top, width, height,
        size=26, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def multiline(slide, lines, left, top, width, height,
              sizes, bolds, colors, aligns, spaces=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aligns[i] if i < len(aligns) else PP_ALIGN.LEFT
        if spaces and i < len(spaces) and spaces[i]:
            p.space_before = Pt(spaces[i])
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sizes[i] if i < len(sizes) else 24)
        run.font.bold = bolds[i] if i < len(bolds) else False
        run.font.color.rgb = colors[i] if i < len(colors) else WHITE
        run.font.name = "Calibri"


def gold_bar(slide, prs, top_inches=1.55):
    W = prs.slide_width
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(top_inches), Inches(12.33), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def bullet_slide_img(prs, blank, title_text, bullets, img_url=None,
                     text_width=0.52):
    """Слайд с заголовком, маркированным списком слева и картинкой справа."""
    s = prs.slides.add_slide(blank)
    set_bg(s, make_bg(), prs)
    W, H = prs.slide_width, prs.slide_height
    gold_bar(s, prs)
    txt(s, title_text, Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.2),
        size=32, bold=True, color=GOLD)
    if img_url:
        add_image_right(s, prs, img_url, left_frac=0.55, top_frac=0.20, w_frac=0.40, h_frac=0.72)
        col_w = Inches(12.33 * text_width)
    else:
        col_w = Inches(12.33)
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.75), col_w, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(5)
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(20 if img_url else 22)
        run.font.color.rgb = LIGHT_GREY
        run.font.name = "Calibri"
    return s


def handler(event: dict, context) -> dict:
    """
    Генерирует PPTX «Космические зонды к внешним планетам: Voyager и открытия колец Урана»
    (Солнечная система). Студент: Ефремова Анастасия Александровна, группа 11С.
    """

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type",
        }, "body": ""}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 1 — ТИТУЛЬНЫЙ
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, make_bg(title=True), prs)

    # Вставка картинки Voyager на правую часть
    add_image_right(s1, prs, IMG_VOYAGER, left_frac=0.53, top_frac=0.12,
                    w_frac=0.43, h_frac=0.76)

    # Надпись-дисциплина
    txt(s1, "АСТРОНОМИЯ  ·  Солнечная система  ·  Тема №22",
        Inches(0.5), Inches(0.35), Inches(7.0), Inches(0.55),
        size=16, color=GOLD, italic=True)

    # Главный заголовок
    multiline(s1,
        ["Космические зонды", "к внешним планетам:", "Voyager и открытия", "колец Урана"],
        Inches(0.5), Inches(1.1), Inches(6.8), Inches(3.5),
        [42, 42, 42, 42], [True, False, False, True],
        [GOLD, WHITE, WHITE, LIGHT_GOLD],
        [PP_ALIGN.LEFT] * 4, [0, 4, 4, 4])

    # Золотая линия-разделитель
    sep = s1.shapes.add_shape(1, Inches(0.5), Inches(5.0), Inches(5.5), Pt(2))
    sep.fill.solid(); sep.fill.fore_color.rgb = GOLD; sep.line.fill.background()

    # Студент
    multiline(s1,
        ["Выполнила: Ефремова Анастасия Александровна", "Группа: 11С"],
        Inches(0.5), Inches(5.2), Inches(6.8), Inches(1.0),
        [20, 20], [False, False], [LIGHT_GREY, LIGHT_GREY],
        [PP_ALIGN.LEFT, PP_ALIGN.LEFT], [0, 5])

    txt(s1, "Исследовательская работа · 2024",
        Inches(0.5), Inches(6.35), Inches(6.8), Inches(0.5),
        size=15, color=GOLD, italic=True)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 2 — АКТУАЛЬНОСТЬ, ЦЕЛЬ, ГИПОТЕЗА
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, make_bg(), prs)
    gold_bar(s2, prs)
    txt(s2, "Актуальность, цель и гипотеза исследования",
        Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.2),
        size=30, bold=True, color=GOLD)

    box2 = s2.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12.33), Inches(5.6))
    tf2 = box2.text_frame; tf2.word_wrap = True
    entries = [
        ("АКТУАЛЬНОСТЬ:",
         "Миссии Voyager (1977–по наст. время) остаются единственными аппаратами, "
         "достигшими Урана и Нептуна. Открытие колец Урана в 1977–1986 гг. пересмотрело "
         "представления о внешних планетах и показало, что кольца — не уникальная черта Сатурна. "
         "Изучение этих данных актуально для подготовки будущих миссий к ледяным гигантам."),
        ("ЦЕЛЬ:",
         "Исследовать историю и результаты миссий Voyager-1 и Voyager-2, "
         "особое внимание уделив открытию и свойствам кольцевой системы Урана, "
         "её отличиям от колец Сатурна и научному значению этих открытий."),
        ("ГИПОТЕЗА:",
         "Кольца Урана образовались позже колец Сатурна и имеют принципиально иное "
         "происхождение — столкновительную природу, а не формирование из протопланетного "
         "диска, — что объясняет их тёмный состав, узость и динамическую нестабильность."),
    ]
    first = True
    for label, body in entries:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        if not first:
            p.space_before = Pt(9)
        r = p.add_run(); r.text = label
        r.font.size = Pt(19); r.font.bold = True
        r.font.color.rgb = GOLD; r.font.name = "Calibri"
        p2 = tf2.add_paragraph(); p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(17); r2.font.color.rgb = LIGHT_GREY; r2.font.name = "Calibri"

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 3 — ПРОГРАММА VOYAGER: ИСТОРИЯ СОЗДАНИЯ
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Программа Voyager: история создания",
        [
            "Запущены в 1977 г. NASA в рамках программы изучения внешних планет.",
            "Воспользовались редким планетным выравниванием (раз в 176 лет): гравитационный манёвр.",
            "Voyager-1 — быстрейший маршрут: Юпитер → Сатурн → за пределы Солнечной системы.",
            "Voyager-2 — «Большое турне»: Юпитер → Сатурн → Уран → Нептун (единственный в истории).",
            "Масса аппарата: ~825 кг; питание — радиоизотопные термоэлектрогенераторы (РТГ).",
            "Связь с Землёй: антенна диаметром 3,7 м, скорость передачи данных 21,6 кбит/с.",
            "На борту — «Золотые пластинки» с посланием человечества внеземным цивилизациям.",
        ],
        IMG_VOYAGER)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 4 — ПУТЬ К УРАНУ: ХРОНОЛОГИЯ ПОЛЁТА
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Путь к Урану: хронология полёта Voyager-2",
        [
            "05.09.1977 — старт Voyager-2 с мыса Канаверал.",
            "1979 г. — пролёт мимо Юпитера, открытие активных вулканов на Ио.",
            "1981 г. — пролёт мимо Сатурна, детальные снимки кольцевой системы.",
            "24.01.1986 — исторический пролёт мимо Урана на расстоянии 81 500 км.",
            "1989 г. — пролёт мимо Нептуна, открытие Большого тёмного пятна.",
            "С 2018 г. — Voyager-2 в межзвёздном пространстве (>18,5 млрд км от Земли).",
            "Сигнал от Voyager-2 идёт до Земли более 17 часов (на 2024 г.).",
        ],
        IMG_SOLAR)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 5 — УРАН: ЛЕДЯНОЙ ГИГАНТ
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Уран: ледяной гигант Солнечной системы",
        [
            "Седьмая планета от Солнца, открыта У. Гершелем в 1781 г.",
            "Диаметр: 50 724 км (в 4 раза больше Земли); масса: 14,5 масс Земли.",
            "Состав: водород, гелий, «ледяная» мантия из воды, аммиака, метана.",
            "Уникальность: ось вращения наклонена на 97,8° — планета «лежит на боку».",
            "Температура атмосферы: −224 °C — самая холодная из планет-гигантов.",
            "27 известных спутников, названных в честь персонажей Шекспира и Поупа.",
            "Voyager-2 открыл 10 новых спутников и два новых кольца Урана.",
        ],
        IMG_URANUS)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 6 — ОТКРЫТИЕ КОЛЕЦ УРАНА
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Открытие колец Урана",
        [
            "10.03.1977 — первое открытие: при покрытии звезды HD 128220 обнаружены 5 колец.",
            "Группа У. Эллиота (MIT): затемнения до и после прохождения Урана — симметричные.",
            "Первоначально обнаружены кольца α, β, γ, δ, ε.",
            "Voyager-2 (1986) уточнил данные и обнаружил ещё 4 кольца: 6, 5, 4 и λ.",
            "Всего каталогизировано 13 известных колец (последние открыты телескопом Хаббл в 2005 г.).",
            "Кольца почти невидимы с Земли: их отражательная способность < 5% (чёрный уголь).",
            "Открытие показало: кольца — не исключение, а норма для планет-гигантов.",
        ],
        IMG_RINGS)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 7 — СОСТАВ И СТРОЕНИЕ КОЛЕЦ УРАНА
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Состав и строение колец Урана",
        [
            "Основной состав: тёмный органический материал (возможно, углеродистые соединения).",
            "Отличие от Сатурна: кольца Урана узкие (шириной 1–100 км), тёмные, не ледяные.",
            "Кольцо ε — самое яркое и широкое (20–96 км), имеет эксцентриситет 0,0079.",
            "Спутники-«пастухи»: Корделия и Офелия удерживают кольцо ε от рассеивания.",
            "Частицы: обломки диаметром от мкм до нескольких метров.",
            "Наклон колец к плоскости орбиты Урана: почти 0° (из-за наклона самой планеты).",
            "Под влиянием солнечного ветра кольца медленно разрушаются (~миллионы лет).",
        ],
        IMG_RINGS)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 8 — НАУЧНЫЕ ОТКРЫТИЯ VOYAGER-2 У УРАНА
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Научные открытия Voyager-2 у Урана",
        [
            "Обнаружено магнитное поле Урана — наклонено на 59° относительно оси вращения.",
            "Открыты 10 новых спутников: Пак, Портия, Джульетта, Крессида, Розалинда и др.",
            "Получены первые детальные снимки Миранды — спутника с уникальным рельефом.",
            "Уточнён период вращения планеты: 17 ч 14 мин.",
            "Обнаружена магнитосфера, асимметричная и смещённая от центра планеты.",
            "Изучено ультрафиолетовое свечение атмосферы — «электросвечение» (electroglow).",
            "Подтверждено: температуры полюсов и экватора Урана практически одинаковы.",
        ],
        IMG_URANUS)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 9 — VOYAGER И НЕПТУН: ПРОДОЛЖЕНИЕ БОЛЬШОГО ТУРНЕ
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Voyager-2 у Нептуна: продолжение Большого турне",
        [
            "25.08.1989 — Voyager-2 достиг Нептуна (3-е место по массе среди планет-гигантов).",
            "Открыто Большое тёмное пятно — атмосферный вихрь размером с Землю.",
            "Обнаружен Тритон — спутник с азотными гейзерами и −235 °C на поверхности.",
            "Открыты кольца Нептуна: Галле, Леверье, Лассель, Арго — узкие и тёмные.",
            "Подтверждено сходство колец Нептуна и Урана: обе системы тёмные и динамически молодые.",
            "Открыто ещё 6 новых спутников Нептуна.",
            "После Нептуна Voyager-2 направился к звёздному пространству.",
        ],
        IMG_NEPTUNE)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 10 — ЗНАЧЕНИЕ МИССИЙ VOYAGER
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide_img(prs, blank,
        "Научное и историческое значение миссий Voyager",
        [
            "Первые и единственные прямые исследования Урана и Нептуна на сегодняшний день.",
            "Открыто более 20 спутников, 2 кольцевые системы, магнитные поля 4 планет.",
            "Voyager-1 (с 2012 г.) и Voyager-2 (с 2018 г.) — первые рукотворные объекты в межзвёздном пространстве.",
            "Данные Voyager стали основой для планирования миссии Uranus Orbiter and Probe (ESA/NASA, 2030-е).",
            "Подтверждена «Гравитационная пращa» — техника экономии топлива при межпланетных полётах.",
            "Доказано: кольцевые системы есть у всех четырёх планет-гигантов Солнечной системы.",
            "Миссии изменили понимание Солнечной системы: внешние планеты оказались разнообразнее, чем ожидалось.",
        ],
        IMG_SOLAR)

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 11 — ВЫВОД: 5 ОГРАНИЧЕНИЙ (ЧТО НЕВОЗМОЖНО)
    # ══════════════════════════════════════════════════════════════════════════
    s11 = prs.slides.add_slide(blank)
    set_bg(s11, make_bg(), prs)
    gold_bar(s11, prs)
    txt(s11, "Вывод: что невозможно — 5 пунктов",
        Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.2),
        size=30, bold=True, color=GOLD)
    box11 = s11.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12.33), Inches(5.6))
    tf11 = box11.text_frame; tf11.word_wrap = True
    negatives = [
        "1. Невозможно провести повторный детальный пролёт мимо Урана до 2030-х гг. — ни одна миссия ещё не отправлена.",
        "2. Невозможно восстановить утраченные данные Voyager-2 о некоторых спутниках Урана: часть памяти была перезаписана.",
        "3. Невозможно точно установить возраст колец Урана: имеющихся данных Voyager-2 недостаточно для однозначного вывода.",
        "4. Невозможно поддерживать прежнюю скорость передачи данных с Voyager: мощность РТГ снизилась на ~70% с 1977 г.",
        "5. Невозможно изменить траекторию Voyager-2: топливо израсходовано, зонд движётся по инерции в межзвёздном пространстве.",
    ]
    for i, line in enumerate(negatives):
        p = tf11.paragraphs[0] if i == 0 else tf11.add_paragraph()
        if i > 0: p.space_before = Pt(9)
        run = p.add_run(); run.text = line
        run.font.size = Pt(19); run.font.color.rgb = RED_SOFT
        run.font.name = "Calibri"

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 12 — ВЫВОД: 5 РЕАЛЬНЫХ ВОЗМОЖНОСТЕЙ
    # ══════════════════════════════════════════════════════════════════════════
    s12 = prs.slides.add_slide(blank)
    set_bg(s12, make_bg(), prs)
    gold_bar(s12, prs)
    txt(s12, "Вывод: реальные возможности — 5 пунктов",
        Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.2),
        size=30, bold=True, color=GOLD)
    box12 = s12.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12.33), Inches(5.6))
    tf12 = box12.text_frame; tf12.word_wrap = True
    positives = [
        "1. Реально: данные Voyager-2 об Уране продолжают анализироваться — в 2023 г. обнаружен плазменный пузырь в магнитосфере.",
        "2. Реально: планируемая миссия Uranus Orbiter (2030-е) позволит изучить кольца и спутники с орбиты многие годы.",
        "3. Реально: применять технику «гравитационной пращи» (открытую Voyager) во всех будущих межпланетных миссиях.",
        "4. Реально: использовать тёмный состав колец Урана как модель для изучения органики во внешней Солнечной системе.",
        "5. Реально: продолжать приём сигналов Voyager-1 и Voyager-2 вплоть до ~2025–2030 гг. до полного истощения РТГ.",
    ]
    for i, line in enumerate(positives):
        p = tf12.paragraphs[0] if i == 0 else tf12.add_paragraph()
        if i > 0: p.space_before = Pt(9)
        run = p.add_run(); run.text = line
        run.font.size = Pt(19); run.font.color.rgb = GREEN_SOFT
        run.font.name = "Calibri"

    # ══════════════════════════════════════════════════════════════════════════
    # Слайд 13 — ПРИЛОЖЕНИЕ: 7 ВОПРОСОВ
    # ══════════════════════════════════════════════════════════════════════════
    s13 = prs.slides.add_slide(blank)
    set_bg(s13, make_bg(), prs)
    gold_bar(s13, prs)
    txt(s13, "Приложение: контрольные вопросы",
        Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.2),
        size=30, bold=True, color=GOLD)
    box13 = s13.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(12.33), Inches(5.6))
    tf13 = box13.text_frame; tf13.word_wrap = True
    questions = [
        "1. Почему 1977 год был выбран для запуска Voyager и какова была его уникальная возможность?",
        "2. Чем отличается маршрут Voyager-1 от маршрута Voyager-2?",
        "3. Как были открыты кольца Урана в 1977 году — каким методом и что наблюдали учёные?",
        "4. Сколько колец Урана было известно до Voyager-2 и сколько обнаружено после его пролёта?",
        "5. Почему кольца Урана так трудно наблюдать с Земли?",
        "6. Какие уникальные особенности Урана (наклон оси, температура, магнитосфера) открыл Voyager-2?",
        "7. Где сейчас находятся Voyager-1 и Voyager-2, и почему их миссия исторически уникальна?",
    ]
    for i, q in enumerate(questions):
        p = tf13.paragraphs[0] if i == 0 else tf13.add_paragraph()
        if i > 0: p.space_before = Pt(8)
        run = p.add_run(); run.text = q
        run.font.size = Pt(19); run.font.color.rgb = LIGHT_GOLD
        run.font.name = "Calibri"

    # ── Сохраняем и возвращаем ────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    data = base64.b64encode(buf.read()).decode("utf-8")

    return {
        "statusCode": 200,
        "headers": {"Access-Control-Allow-Origin": "*", "Content-Type": "application/json"},
        "body": '{"filename":"voyager_uranus_rings.pptx","data":"' + data + '"}',
    }